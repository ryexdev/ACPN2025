import streamlit as st
import pandas as pd
import os
import io
import base64
from pathlib import Path
from classes.ai_engines.openai_client import openai_client
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import seaborn as sns

# Page configuration
st.set_page_config(
    page_title="KPI Analyzer",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Get API key from environment variable or ask user
secret_value = os.getenv("OwadmasdujU")
model_name = "gpt-4.1-nano"

# Title and description
st.title("📊 KPI Analyzer")
with st.expander("Description of KPI Analyzer", expanded=True):
    st.markdown("""
    Transform your raw data into insightful KPI reports and presentations. This powerful tool helps you:

    - Upload and analyze CSV, TSV, or Excel files containing your business metrics
    - Generate automated insights and trend analysis using AI
    - Create professional PowerPoint presentations with data visualizations
    - Compare performance against targets and identify key trends
    - Export ready-to-use reports for meetings and presentations
    - Use sample datasets to explore the tool's capabilities

    Whether you're analyzing sales performance, customer satisfaction, or inventory metrics, 
    this tool helps you turn complex data into clear, actionable insights.
    """)

# Get API key from environment variable or ask user to enter it
if not secret_value:
    with st.container(border=True):
        st.subheader("OpenAI API Key")
        st.warning("API key is not set. Please enter your API key below to continue to use the tool.")
        api_key = st.text_input("Enter your API key:", type="password")
        model_name = st.selectbox(
            "Select OpenAI Model:", 
            ["gpt-4.1-nano", "gpt-4o-mini"],
            index=0
        )
else:
    st.success("OpenAI API key has been provided for the demo. You can freely use the tool until the API key expires (estimated 2025-05-14 @ 12:00 MST).")
    api_key = secret_value

st.divider()

# Check if required libraries are installed
try:
    import python_pptx
except ImportError:
    st.info("Installing required libraries...")
    os.system("pip install python-pptx matplotlib seaborn")
    st.success("Libraries installed!")
except Exception as e:
    st.error(f"Error installing libraries: {e}. Contact the developer for assistance.")
    st.stop()

# File upload section
st.header("1. Upload Your Data")
uploaded_file = st.file_uploader(
    "Upload CSV, TSV, or Excel file", 
    type=["csv", "tsv", "xlsx", "xls"], 
    help="Upload your data file to analyze"
)

# Sample data section
with st.expander("Or use sample data"):
    sample_option = st.radio(
        "Select sample data:",
        ["Sales Performance", "Customer Satisfaction", "Inventory Turnover"]
    )
    
    if st.button("Load Sample Data"):
        if sample_option == "Sales Performance":
            # Create sample sales data
            data = {
                "Month": ["Jan", "Feb", "Mar", "Apr", "May", "Jun", 
                         "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
                "Region": ["East"] * 4 + ["West"] * 4 + ["North"] * 4,
                "Sales": [12500, 14200, 15100, 16400, 11800, 13200, 
                         14500, 15800, 10200, 11500, 12800, 14100],
                "Target": [13000, 13000, 14000, 15000, 12000, 13000, 
                         14000, 15000, 11000, 12000, 13000, 14000],
                "Growth": [0, 13.6, 6.3, 8.6, 0, 11.9, 9.8, 9.0, 0, 12.7, 11.3, 10.2]
            }
            df = pd.DataFrame(data)
            st.session_state['df'] = df
            st.session_state['file_loaded'] = True
            st.success("Sample sales data loaded successfully!")
            
        elif sample_option == "Customer Satisfaction":
            # Create sample customer satisfaction data
            data = {
                "Quarter": ["Q1", "Q1", "Q1", "Q2", "Q2", "Q2", "Q3", "Q3", "Q3", "Q4", "Q4", "Q4"],
                "Department": ["Sales", "Support", "Product"] * 4,
                "Satisfaction": [85, 78, 82, 87, 80, 84, 89, 82, 86, 91, 83, 88],
                "Response Rate": [65, 72, 58, 68, 75, 62, 71, 78, 64, 73, 80, 67],
                "NPS": [42, 35, 38, 45, 38, 40, 48, 41, 43, 50, 43, 46]
            }
            df = pd.DataFrame(data)
            st.session_state['df'] = df
            st.session_state['file_loaded'] = True
            st.success("Sample customer satisfaction data loaded successfully!")
            
        elif sample_option == "Inventory Turnover":
            # Create sample inventory data
            data = {
                "Product_Category": ["Brakes", "Filters", "Lighting", "Engine", "Suspension", 
                                    "Electrical", "Body", "Fluids", "Cooling", "Interior"],
                "Units_Sold": [325, 512, 189, 276, 198, 423, 156, 387, 209, 143],
                "Inventory_Level": [78, 125, 56, 83, 62, 91, 42, 103, 58, 39],
                "Restock_Time_Days": [5, 3, 7, 10, 8, 4, 6, 2, 9, 11],
                "Profit_Margin": [32, 28, 35, 41, 38, 29, 36, 25, 33, 39]
            }
            df = pd.DataFrame(data)
            st.session_state['df'] = df
            st.session_state['file_loaded'] = True
            st.success("Sample inventory data loaded successfully!")

# Process uploaded file
if uploaded_file is not None:
    try:
        # Determine file type and read accordingly
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        if file_extension in ['xlsx', 'xls']:
            df = pd.read_excel(uploaded_file)
        elif file_extension == 'csv':
            df = pd.read_csv(uploaded_file)
        elif file_extension == 'tsv':
            df = pd.read_csv(uploaded_file, sep='\t')
        
        st.session_state['df'] = df
        st.session_state['file_loaded'] = True
        st.success(f"File '{uploaded_file.name}' loaded successfully!")
    except Exception as e:
        st.error(f"Error loading file: {e}")

# Display data if loaded
if 'file_loaded' in st.session_state and st.session_state['file_loaded']:
    df = st.session_state['df']
    
    st.header("2. Review Your Data")
    
    # Basic data info
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Data Preview")
        st.dataframe(df.head(10), use_container_width=True)
    
    with col2:
        st.subheader("Data Summary")
        st.write(f"**Rows:** {df.shape[0]}")
        st.write(f"**Columns:** {df.shape[1]}")
        st.write("**Column Types:**")
        st.write(df.dtypes)
    
    # KPI Analysis Settings
    st.header("3. KPI Analysis Settings")
    
    # Select columns for analysis
    numeric_cols = df.select_dtypes(include=['float64', 'int64']).columns.tolist()
    categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("Select Metrics")
        selected_metrics = st.multiselect(
            "Choose metrics for analysis:",
            numeric_cols,
            default=numeric_cols[:min(3, len(numeric_cols))]
        )
    
    with col2:
        st.subheader("Select Dimensions")
        selected_dimensions = st.multiselect(
            "Choose dimensions for grouping:",
            categorical_cols,
            default=categorical_cols[:min(1, len(categorical_cols))]
        )
    
    # OpenAI API settings
    st.header("4. AI Analysis Settings")
    
        
    # Analysis focus
    analysis_focus = st.selectbox(
        "Analysis Focus:",
        ["General Performance Overview", 
            "Trend Analysis and Forecasting", 
            "Key Improvement Areas", 
            "Comparative Analysis", 
            "Executive Summary"],
        index=0
    )
    
    # Create the analysis and presentation
    if st.button("Generate KPI Analysis & Presentation", type="primary"):
        if not selected_metrics:
            st.warning("Please select at least one metric for analysis.")
        elif not api_key and not secret_value:
            st.warning("Please enter your OpenAI API key.")
        else:
            with st.spinner("Analyzing data and creating presentation..."):
                try:
                    # Create a progress container
                    progress_container = st.container()
                    progress_container.info("Starting analysis...")
                    
                    # Create filtered dataframe with selected columns
                    analysis_cols = selected_dimensions + selected_metrics
                    analysis_df = df[analysis_cols].copy()
                    
                    # Generate chart images
                    charts = []
                    
                    progress_container.info("Generating charts...")
                    for metric in selected_metrics:
                        plt.figure(figsize=(10, 6))
                        
                        if len(selected_dimensions) > 0:
                            dimension = selected_dimensions[0]
                            
                            # Generate appropriate chart based on data types
                            if len(df[dimension].unique()) <= 10:  # For fewer categories
                                ax = sns.barplot(x=dimension, y=metric, data=df)
                                plt.title(f"{metric} by {dimension}")
                                plt.xticks(rotation=45)
                                plt.tight_layout()
                            else:  # For many categories
                                ax = sns.lineplot(x=dimension, y=metric, data=df)
                                plt.title(f"{metric} Trend by {dimension}")
                                plt.xticks(rotation=45)
                                plt.tight_layout()
                        else:
                            # If no dimensions, create histogram or box plot
                            fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 6))
                            sns.histplot(df[metric], kde=True, ax=ax1)
                            ax1.set_title(f"Distribution of {metric}")
                            
                            sns.boxplot(y=df[metric], ax=ax2)
                            ax2.set_title(f"Box Plot of {metric}")
                            plt.tight_layout()
                        
                        # Save chart to memory
                        chart_buffer = io.BytesIO()
                        plt.savefig(chart_buffer, format='png', dpi=100)
                        chart_buffer.seek(0)
                        charts.append(chart_buffer)
                        plt.close()
                    
                    # Create AI prompt
                    progress_container.info("Requesting AI analysis...")
                    data_sample = df.head(20).to_string()
                    data_summary = df.describe().to_string()
                    
                    if len(selected_dimensions) > 0:
                        dimension_summary = df.groupby(selected_dimensions[0])[selected_metrics].mean().to_string()
                    else:
                        dimension_summary = "No dimensions selected for grouping."
                    
                    prompt = f"""
                    As a business intelligence analyst, analyze the following data and create a KPI report focused on {analysis_focus}.
                    
                    DATA SAMPLE:
                    {data_sample}
                    
                    DATA SUMMARY:
                    {data_summary}
                    
                    DIMENSION SUMMARY:
                    {dimension_summary}
                    
                    Please provide:
                    1. An executive summary (3-4 sentences)
                    2. Key insights (5 bullet points)
                    3. Recommendations (3 bullet points)
                    4. Analysis of each metric: {', '.join(selected_metrics)} (2-3 sentences per metric)
                    
                    Format your response as a structured report with these clear sections.
                    """
                    
                    # Generate AI analysis
                    ai_analysis = openai_client.generate_with_openai(prompt, model_name)
                    
                    # Create PowerPoint presentation
                    progress_container.info("Creating PowerPoint presentation...")
                    ppt = Presentation()
                    
                    # Title slide
                    title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
                    title = title_slide.shapes.title
                    subtitle = title_slide.placeholders[1]
                    title.text = "KPI Analysis Report"
                    subtitle.text = f"Focus: {analysis_focus}"
                    
                    # Summary slide
                    summary_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                    summary_slide.shapes.title.text = "Executive Summary"
                    summary_text = summary_slide.placeholders[1]
                    
                    # Extract executive summary from AI analysis
                    exec_summary = ai_analysis.split("Key insights")[0].strip()
                    if "Executive summary" in exec_summary:
                        exec_summary = exec_summary.split("Executive summary")[1].strip()
                    
                    summary_text.text = exec_summary
                    
                    # Key insights slide
                    insights_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                    insights_slide.shapes.title.text = "Key Insights"
                    insights_text = insights_slide.placeholders[1]
                    
                    # Extract key insights from AI analysis
                    key_insights = ""
                    if "Key insights" in ai_analysis and "Recommendations" in ai_analysis:
                        key_insights = ai_analysis.split("Key insights")[1].split("Recommendations")[0].strip()
                    
                    insights_text.text = key_insights
                    
                    # Charts slides
                    for i, chart_buffer in enumerate(charts):
                        chart_slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                        chart_slide.shapes.title.text = f"Analysis: {selected_metrics[i]}"
                        
                        # Add chart image
                        chart_buffer.seek(0)
                        chart_data = chart_buffer.getvalue()
                        pic = chart_slide.shapes.add_picture(io.BytesIO(chart_data), Inches(2), Inches(2), width=Inches(6))
                        
                        # Add AI analysis for this metric
                        txbox = chart_slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
                        text_frame = txbox.text_frame
                        text_frame.text = f"Analysis: "
                        
                        # Find specific analysis for this metric
                        metric_name = selected_metrics[i]
                        metric_analysis = ""
                        if f"Analysis of {metric_name}" in ai_analysis:
                            metric_analysis = ai_analysis.split(f"Analysis of {metric_name}")[1].split("Analysis of ")[0].strip()
                        elif metric_name in ai_analysis:
                            parts = ai_analysis.split(metric_name)
                            if len(parts) > 1:
                                for part in parts[1:]:
                                    if len(part.strip()) > 0 and len(part.strip()) < 500:  # Reasonable length for analysis
                                        metric_analysis = part.strip()
                                        break
                        
                        if not metric_analysis:
                            metric_analysis = "See overall analysis in recommendations slide."
                            
                        p = text_frame.add_paragraph()
                        p.text = metric_analysis
                    
                    # Recommendations slide
                    recommendations_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                    recommendations_slide.shapes.title.text = "Recommendations"
                    recommendations_text = recommendations_slide.placeholders[1]
                    
                    # Extract recommendations from AI analysis
                    recommendations = ""
                    if "Recommendations" in ai_analysis:
                        recommendations = ai_analysis.split("Recommendations")[1].split("Analysis of ")[0].strip()
                    
                    recommendations_text.text = recommendations
                    
                    # Save presentation
                    ppt_buffer = io.BytesIO()
                    ppt.save(ppt_buffer)
                    ppt_buffer.seek(0)
                    
                    # Success message
                    progress_container.success("Analysis and presentation complete!")
                    
                    # Display AI analysis with collapsible section
                    st.header("5. AI Analysis Results")
                    
                    with st.expander("View AI Analysis Details", expanded=True):
                        st.subheader("OpenAI API Response")
                        st.code(f"Model used: {model_name}\n\nPrompt sent to OpenAI:\n{prompt}\n\nResponse:\n{ai_analysis}")
                    
                    # Display a sample chart
                    st.subheader("Sample Visualization")
                    if charts:
                        charts[0].seek(0)
                        st.image(charts[0], use_column_width=True)
                    
                    # Download PowerPoint
                    st.header("6. Download Presentation")
                    
                    # Create download button for PPT
                    st.download_button(
                        "Download PowerPoint Presentation",
                        ppt_buffer.getvalue(),
                        "KPI_Analysis.pptx",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key='download-ppt'
                    )
                    
                except Exception as e:
                    st.error(f"Error generating analysis and presentation: {e}")

# Create sample data files for download
st.header("Sample Data Files")
st.markdown("Download these sample files to test the KPI Analyzer:")

# Create sample Excel file
excel_buffer = io.BytesIO()
with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
    # Sales data
    sales_data = pd.DataFrame({
        "Month": ["Jan", "Feb", "Mar", "Apr", "May", "Jun", 
                 "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
        "Region": ["East"] * 4 + ["West"] * 4 + ["North"] * 4,
        "Sales": [12500, 14200, 15100, 16400, 11800, 13200, 
                 14500, 15800, 10200, 11500, 12800, 14100],
        "Target": [13000, 13000, 14000, 15000, 12000, 13000, 
                 14000, 15000, 11000, 12000, 13000, 14000],
        "Growth": [0, 13.6, 6.3, 8.6, 0, 11.9, 9.8, 9.0, 0, 12.7, 11.3, 10.2]
    })
    sales_data.to_excel(writer, sheet_name='Sales Performance', index=False)
    
    # Customer satisfaction data
    csat_data = pd.DataFrame({
        "Quarter": ["Q1", "Q1", "Q1", "Q2", "Q2", "Q2", "Q3", "Q3", "Q3", "Q4", "Q4", "Q4"],
        "Department": ["Sales", "Support", "Product"] * 4,
        "Satisfaction": [85, 78, 82, 87, 80, 84, 89, 82, 86, 91, 83, 88],
        "Response Rate": [65, 72, 58, 68, 75, 62, 71, 78, 64, 73, 80, 67],
        "NPS": [42, 35, 38, 45, 38, 40, 48, 41, 43, 50, 43, 46]
    })
    csat_data.to_excel(writer, sheet_name='Customer Satisfaction', index=False)
    
    # Inventory data
    inventory_data = pd.DataFrame({
        "Product_Category": ["Brakes", "Filters", "Lighting", "Engine", "Suspension", 
                            "Electrical", "Body", "Fluids", "Cooling", "Interior"],
        "Units_Sold": [325, 512, 189, 276, 198, 423, 156, 387, 209, 143],
        "Inventory_Level": [78, 125, 56, 83, 62, 91, 42, 103, 58, 39],
        "Restock_Time_Days": [5, 3, 7, 10, 8, 4, 6, 2, 9, 11],
        "Profit_Margin": [32, 28, 35, 41, 38, 29, 36, 25, 33, 39]
    })
    inventory_data.to_excel(writer, sheet_name='Inventory Turnover', index=False)

excel_buffer.seek(0)

# Create sample TSV file
tsv_buffer = io.StringIO()
sales_data.to_csv(tsv_buffer, sep='\t', index=False)
tsv_content = tsv_buffer.getvalue()
tsv_bytes = tsv_content.encode()

# Download buttons for sample files
col1, col2 = st.columns(2)

with col1:
    st.download_button(
        "Download Sample Excel File",
        excel_buffer,
        "Sample_KPI_Data.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        key='download-xlsx'
    )

with col2:
    st.download_button(
        "Download Sample TSV File",
        tsv_bytes,
        "Sample_Sales_Data.tsv",
        "text/tab-separated-values",
        key='download-tsv'
    )

# Help and Information
with st.expander("How to Use This Tool"):
    st.markdown("""
    ### 📊 KPI Analyzer - User Guide
    
    This tool helps you analyze your business data and create professional KPI presentations with AI assistance.
    
    #### Steps to use:
    
    1. **Upload Data**: Upload your CSV, TSV, or Excel file containing your KPI data.
    2. **Review Data**: Check the preview to ensure your data was loaded correctly.
    3. **Configure Analysis**: Select the metrics and dimensions you want to analyze.
    4. **Set AI Parameters**: Enter your OpenAI API key and select analysis focus.
    5. **Generate Analysis**: Click the button to create your KPI analysis and presentation.
    6. **Download Results**: Get your PowerPoint presentation for immediate use.
    
    #### Tips for best results:
    
    - Ensure your data has clear column names
    - Include both numeric metrics and categorical dimensions
    - For time-series data, make sure dates are in chronological order
    - The more data you provide, the better the AI analysis will be
    - Select 2-4 key metrics for focused insights
    """)

# Footer
st.markdown("---")
st.markdown("✨ *Powered by OpenAI and Streamlit* ✨")
